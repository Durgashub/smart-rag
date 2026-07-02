"""
Ingest documents for a specific user session.

Usage:
    python ingest.py --session <session_id>

Reads from:  docs/<session_id>/
Writes to:   vector_store/<session_id>/

Chunking strategy: Parent-Child hierarchy
- Small child chunks (400 chars) used for precise FAISS retrieval
- Large parent chunks (1200 chars) stored alongside for GPT context
- At retrieval time, child chunk finds the match, parent chunk is sent to GPT

Supported file types:
    .pdf   — pdfplumber + OCR fallback
    .docx  — python-docx
    .xlsx  — openpyxl
    .csv   — csv module
    .pptx  — python-pptx
    .txt / .md — plain text
"""

import os
import re
import json
import glob
import argparse
from pathlib import Path

import numpy as np
import faiss
from dotenv import load_dotenv
from openai import OpenAI
import pdfplumber
from pdf2image import convert_from_path
import pytesseract

load_dotenv()
client = OpenAI()

EMBED_MODEL = "text-embedding-3-small"

# Parent-child chunk sizes
PARENT_CHUNK_SIZE = 1200   # sent to GPT for answering
CHILD_CHUNK_SIZE = 400     # used for retrieval (more precise)
CHILD_OVERLAP = 50         # small overlap between child chunks

CHUNK_OVERLAP = 200        # overlap between parent chunks
MIN_TEXT_LEN_BEFORE_OCR = 20

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv",
    ".txt", ".md", ".pptx"
}

TESSERACT_CMD = os.getenv("TESSERACT_CMD")
if TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
POPPLER_PATH = os.getenv("POPPLER_PATH") or None


# ── Extractors ────────────────────────────────────────────────────────────────

def _table_to_markdown(table_rows) -> str:
    rows = [[(cell or "").strip().replace("\n", " ") for cell in row] for row in table_rows if row]
    if not rows:
        return ""
    header = rows[0]
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        row = (row + [""] * len(header))[: len(header)]
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _ocr_page(pdf_path: str, page_number: int) -> str:
    images = convert_from_path(
        pdf_path, first_page=page_number, last_page=page_number,
        dpi=200, poppler_path=POPPLER_PATH,
    )
    if not images:
        return ""
    return pytesseract.image_to_string(images[0])


def load_text_from_pdf(path: str) -> str:
    page_texts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            tables = page.extract_tables()
            table_blocks = [_table_to_markdown(t) for t in tables if t]
            table_blocks = [t for t in table_blocks if t]
            combined = text
            if table_blocks:
                combined = (combined + "\n\n" + "\n\n".join(table_blocks)).strip()
            if len(combined) < MIN_TEXT_LEN_BEFORE_OCR:
                try:
                    ocr_text = _ocr_page(path, i)
                    if ocr_text.strip():
                        combined = ocr_text.strip()
                except Exception as e:
                    print(f"  Warning: OCR failed on page {i}: {e}")
            if combined:
                page_texts.append(f"[Page {i}]\n{combined}")
    return "\n\n".join(page_texts)


def load_text_from_docx(path: str) -> str:
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        for table in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if rows:
                parts.append(_table_to_markdown(rows))
        return "\n\n".join(parts)
    except Exception as e:
        print(f"  Warning: Could not read docx {Path(path).name}: {e}")
        return ""


def load_text_from_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_data):
                    rows.append(row_data)
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + _table_to_markdown(rows))
        return "\n\n".join(parts)
    except ImportError:
        print("  Warning: openpyxl not installed.")
        return ""
    except Exception as e:
        print(f"  Warning: Could not read xlsx: {e}")
        return ""


def load_text_from_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if slide_texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except ImportError:
        print("  Warning: python-pptx not installed.")
        return ""
    except Exception as e:
        print(f"  Warning: Could not read pptx: {e}")
        return ""


def load_text_from_csv(path: str) -> str:
    try:
        import csv
        rows = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(row)
        return _table_to_markdown(rows) if rows else ""
    except Exception as e:
        print(f"  Warning: Could not read csv: {e}")
        return ""


def load_text_from_file(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return load_text_from_pdf(path)
    if ext == ".docx":
        return load_text_from_docx(path)
    if ext == ".doc":
        print(f"  Warning: .doc not supported directly. Convert to .docx.")
        return ""
    if ext in {".xlsx", ".xls"}:
        return load_text_from_xlsx(path)
    if ext == ".csv":
        return load_text_from_csv(path)
    if ext == ".pptx":
        return load_text_from_pptx(path)
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        if content.count("\x00") > 10:
            print(f"  Warning: {Path(path).name} appears binary, skipping.")
            return ""
        return content
    except Exception as e:
        print(f"  Warning: Could not read {Path(path).name}: {e}")
        return ""


# ── Low-value chunk filter ────────────────────────────────────────────────────

def is_low_value_chunk(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return True
    lines = [l for l in stripped.splitlines() if l.strip()]
    if not lines:
        return True
    table_lines = sum(1 for l in lines if l.strip().startswith("|"))
    if table_lines == len(lines) and "CONCEPT CHECK" in stripped.upper():
        return True
    numbered = sum(1 for l in lines if re.match(r"^\d{1,3}[\.\)]\s", l.strip()))
    if len(lines) >= 3 and numbered / len(lines) > 0.5:
        return True
    return False


# ── Parent-child chunking ─────────────────────────────────────────────────────

def chunk_into_parents(
    text: str,
    chunk_size: int = PARENT_CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> list[str]:
    """
    Split text into large parent chunks (paragraph-aware).
    These are stored for GPT context — bigger = more complete answers.
    """
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks = []
    current = ""

    def flush():
        nonlocal current
        if current.strip():
            chunks.append(current.strip())
        current = ""

    for para in paragraphs:
        if len(para) > chunk_size:
            flush()
            start = 0
            while start < len(para):
                chunks.append(para[start : start + chunk_size].strip())
                start += chunk_size - overlap
            continue
        if len(current) + len(para) + 2 <= chunk_size:
            current = f"{current}\n\n{para}" if current else para
        else:
            flush()
            current = para

    flush()
    return [c for c in chunks if c]


def split_into_children(
    parent: str,
    child_size: int = CHILD_CHUNK_SIZE,
    overlap: int = CHILD_OVERLAP,
) -> list[str]:
    """
    Split a parent chunk into smaller child chunks.
    These are used for retrieval — smaller = more precise matching.
    """
    children = []
    start = 0
    while start < len(parent):
        end = start + child_size
        children.append(parent[start:end].strip())
        start += child_size - overlap
    return [c for c in children if len(c) > 50]  # skip tiny fragments


def build_parent_child_chunks(text: str) -> list[dict]:
    """
    Build parent-child chunk pairs from document text.

    Returns list of dicts:
    {
        "child_text":  "...400 char precise chunk for retrieval...",
        "parent_text": "...1200 char full chunk sent to GPT...",
        "parent_id":   0,  # which parent this child belongs to
    }

    At retrieval time:
    - FAISS indexes child_text embeddings (precise matching)
    - GPT receives parent_text (full context for better answers)
    """
    parents = chunk_into_parents(text)
    result = []

    for parent_id, parent in enumerate(parents):
        if is_low_value_chunk(parent):
            continue

        children = split_into_children(parent)

        for child in children:
            if is_low_value_chunk(child):
                continue
            result.append({
                "child_text": child,
                "parent_text": parent,
                "parent_id": parent_id,
            })

    return result


# ── Embeddings ────────────────────────────────────────────────────────────────

def embed_texts(texts: list[str]) -> list[list[float]]:
    if not texts:
        return []
    all_embeddings = []
    batch_size = 100
    for i in range(0, len(texts), batch_size):
        batch = texts[i : i + batch_size]
        response = client.embeddings.create(model=EMBED_MODEL, input=batch)
        all_embeddings.extend(item.embedding for item in response.data)
    return all_embeddings


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--session", required=True)
    args = parser.parse_args()

    session_id = args.session
    docs_dir = f"docs/{session_id}"
    store_dir = f"vector_store/{session_id}"

    os.makedirs(store_dir, exist_ok=True)

    all_files = glob.glob(f"{docs_dir}/**/*.*", recursive=True)
    file_paths = [
        p for p in all_files
        if Path(p).suffix.lower() in SUPPORTED_EXTENSIONS
    ]

    if not file_paths:
        print(f"No supported files found in '{docs_dir}/'.")
        return

    # metadata now stores both child and parent text
    all_child_texts = []   # embedded for retrieval
    metadata = []          # stored alongside index
    total_parents = 0
    total_children = 0
    skipped = 0

    for path in file_paths:
        ext = Path(path).suffix.lower()
        print(f"Loading {path} ({ext}) ...")
        text = load_text_from_file(path)

        if not text.strip():
            print(f"  -> No text extracted, skipping.")
            continue

        pairs = build_parent_child_chunks(text)

        if not pairs:
            print(f"  -> No chunks produced, skipping.")
            continue

        # Count unique parents
        unique_parents = len(set(p["parent_id"] for p in pairs))
        total_parents += unique_parents
        total_children += len(pairs)
        skipped_raw = len(chunk_into_parents(text)) - unique_parents
        skipped += skipped_raw

        print(
            f"  -> {unique_parents} parent chunks → "
            f"{len(pairs)} child chunks "
            f"({skipped_raw} low-value parents skipped)"
        )

        for pair in pairs:
            all_child_texts.append(pair["child_text"])
            metadata.append({
                "source": Path(path).name,
                # text = parent_text: GPT gets full context
                "text": pair["parent_text"],
                # child_text stored for debug/display
                "child_text": pair["child_text"],
                "parent_id": pair["parent_id"],
            })

    if not all_child_texts:
        print("No child chunks produced. Check file contents.")
        return

    print(f"\nTotal: {total_parents} parents → {total_children} children across {len(file_paths)} file(s)")
    print(f"Embedding {len(all_child_texts)} child chunks (used for retrieval)...")

    # Embed child chunks for precise retrieval
    embeddings = embed_texts(all_child_texts)

    embedding_matrix = np.array(embeddings, dtype="float32")
    dimension = embedding_matrix.shape[1]

    index = faiss.IndexFlatL2(dimension)
    index.add(embedding_matrix)

    faiss.write_index(index, f"{store_dir}/index.faiss")
    with open(f"{store_dir}/metadata.json", "w", encoding="utf-8") as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)

    print(f"Done. Index saved to '{store_dir}/'.")
    print(f"  Child chunks indexed: {len(all_child_texts)}")
    print(f"  GPT will receive parent chunks ({PARENT_CHUNK_SIZE} chars) for richer answers.")


if __name__ == "__main__":
    main()