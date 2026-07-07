"""
ingestion/loaders.py — document text extraction for all supported formats.

To add a new format:
  1. Write a load_text_from_X() function
  2. Register it in load_text_from_file()
"""

import os
from pathlib import Path
import pytesseract
from config import TESSERACT_CMD, POPPLER_PATH, MIN_TEXT_LEN_BEFORE_OCR

if TESSERACT_CMD:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD


def _table_to_markdown(table_rows) -> str:
    """Convert a pdfplumber table (list of row lists) to Markdown grid."""
    rows = [
        [(cell or "").strip().replace("\n", " ") for cell in row]
        for row in table_rows if row
    ]
    if not rows:
        return ""
    header = rows[0]
    lines  = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in rows[1:]:
        row = (row + [""] * len(header))[:len(header)]
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _ocr_page(pdf_path: str, page_number: int) -> str:
    from pdf2image import convert_from_path
    images = convert_from_path(
        pdf_path, first_page=page_number, last_page=page_number,
        dpi=200, poppler_path=POPPLER_PATH,
    )
    if not images:
        return ""
    return pytesseract.image_to_string(images[0])


def load_text_from_pdf(path: str) -> str:
    import pdfplumber
    page_texts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text        = (page.extract_text() or "").strip()
            tables      = page.extract_tables()
            table_blocks = [_table_to_markdown(t) for t in tables if t]
            table_blocks = [t for t in table_blocks if t]
            combined    = text
            if table_blocks:
                combined = (combined + "\n\n" + "\n\n".join(table_blocks)).strip()
            if len(combined) < MIN_TEXT_LEN_BEFORE_OCR:
                try:
                    ocr_text = _ocr_page(path, i)
                    if ocr_text.strip():
                        combined = ocr_text.strip()
                except Exception as e:
                    print(f"  Warning: OCR failed on page {i}: {e}")
            if combined:
                page_texts.append(f"[Page {i}]\n{combined}")
    return "\n\n".join(page_texts)


def load_text_from_docx(path: str) -> str:
    try:
        from docx import Document
        doc   = Document(path)
        parts = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if rows:
                parts.append(_table_to_markdown(rows))
        return "\n\n".join(parts)
    except Exception as e:
        print(f"  Warning: Could not read docx {Path(path).name}: {e}")
        return ""


def load_text_from_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb    = openpyxl.load_workbook(path, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws   = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                if any(cell.strip() for cell in row_data):
                    rows.append(row_data)
            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + _table_to_markdown(rows))
        return "\n\n".join(parts)
    except ImportError:
        print("  Warning: openpyxl not installed.")
        return ""
    except Exception as e:
        print(f"  Warning: Could not read xlsx: {e}")
        return ""


def load_text_from_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(parts)
    except ImportError:
        print("  Warning: python-pptx not installed.")
        return ""
    except Exception as e:
        print(f"  Warning: Could not read pptx: {e}")
        return ""


def load_text_from_csv(path: str) -> str:
    try:
        import csv
        rows = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            for row in csv.reader(f):
                if any(cell.strip() for cell in row):
                    rows.append(row)
        return _table_to_markdown(rows) if rows else ""
    except Exception as e:
        print(f"  Warning: Could not read csv: {e}")
        return ""


def load_text_from_file(path: str) -> str:
    """Route a file to the correct loader based on extension."""
    ext = Path(path).suffix.lower()
    loaders = {
        ".pdf":  load_text_from_pdf,
        ".docx": load_text_from_docx,
        ".xlsx": load_text_from_xlsx,
        ".xls":  load_text_from_xlsx,
        ".csv":  load_text_from_csv,
        ".pptx": load_text_from_pptx,
    }
    loader = loaders.get(ext)
    if loader:
        return loader(path)
    # Plain text fallback
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        if content.count("\x00") > 10:
            print(f"  Warning: {Path(path).name} appears binary, skipping.")
            return ""
        return content
    except Exception as e:
        print(f"  Warning: Could not read {Path(path).name}: {e}")
        return ""
